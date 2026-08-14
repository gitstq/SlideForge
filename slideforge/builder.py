"""原生 .pptx 渲染器。

基于 python-pptx 把结构化大纲渲染为原生 PowerPoint 文件：
16:9 宽屏、主题配色、原生形状、原生图表、演讲备注、
原生切换动画与逐页入场动画（通过 OOXML 注入，保持 100% 兼容）。
"""

from __future__ import annotations

import io
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .outline import parse_outline, sample_outline
from .themes import Theme, get_theme

W, H = 13.333, 7.5  # 16:9 英寸

# 版式对应的入场动画（presetID，0 = 无动画）
ENTRANCE_PRESETS = {"fade": 10, "rise": 25, "zoom": 23, "none": 0}


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


class SlideBuilder:
    def __init__(self, theme: Theme, effect: str = "fade", transition: str = "fade"):
        self.theme = theme
        self.effect = effect if effect in ENTRANCE_PRESETS else "fade"
        self.transition = transition
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self._blank = self.prs.slide_layouts[6]

    # ---------- 基础工具 ----------

    def _new_slide(self, accent: Optional[str] = None) -> "pptx.slide.Slide":
        slide = self.prs.slides.add_slide(self._blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(self.theme.bg)
        self._add_transition(slide, self.transition)
        self._remember_accent(accent or self.theme.accent)
        return slide

    def _remember_accent(self, accent: str):
        # 供后续形状复用当前页强调色
        self._cur_accent = accent

    @property
    def accent(self) -> str:
        return getattr(self, "_cur_accent", self.theme.accent)

    def _textbox(self, slide, x, y, w, h, text="", size=18, color=None,
                 bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font=None, line_spacing=None, wrap=True, valign_center=True):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        p.space_after = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = _rgb(color or self.theme.ink)
        f.name = font or self.theme.heading_font
        return tb, tf

    def _shape(self, slide, stype, x, y, w, h, fill=None, line=None):
        shp = slide.shapes.add_shape(stype, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(fill)
        else:
            shp.fill.background()
        if line:
            shp.line.color.rgb = _rgb(line)
            shp.line.width = Pt(1.2)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _set_text_in_shape(self, shp, text, size=16, color=None, bold=False,
                           align=PP_ALIGN.CENTER, font=None, anchor=MSO_ANCHOR.MIDDLE):
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top = tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = _rgb(color or self.theme.ink)
        f.name = font or self.theme.body_font

    def _bullets_frame(self, slide, x, y, w, h, bullets, size=17,
                       marker="square", color=None, spacing=1.15):
        tb, tf = self._textbox(slide, x, y, w, h, size=size,
                               color=color or self.theme.ink)
        first = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(size * 0.5)
            p.line_spacing = spacing
            # marker run
            if marker == "square":
                mr = p.add_run()
                mr.text = "\u25A0  "
                mf = mr.font
                mf.size = Pt(max(10, int(size * 0.55)))
                mf.color.rgb = _rgb(self.accent)
                mf.name = self.theme.body_font
            r = p.add_run()
            r.text = str(b)
            rf = r.font
            rf.size = Pt(size)
            rf.color.rgb = _rgb(color or self.theme.ink)
            rf.name = self.theme.body_font
        return tb

    # ---------- 原生动画 / 过渡 ----------

    def _add_transition(self, slide, kind: str):
        try:
            sld = slide._element
            for el in sld.findall(qn("p:transition")):
                sld.remove(el)
            trans = sld.makeelement(qn("p:transition"), {"spd": "med"})
            t = kind or "fade"
            if t == "random":
                trans.append(sld.makeelement(qn("p:random"), {}))
            elif t == "push":
                trans.append(sld.makeelement(qn("p:push"), {"dir": "l"}))
            elif t == "split":
                trans.append(sld.makeelement(qn("p:split"), {"orient": "horz", "dir": "out"}))
            elif t == "wipe":
                trans.append(sld.makeelement(qn("p:wipe"), {"dir": "l"}))
            else:
                trans.append(sld.makeelement(qn("p:fade"), {}))
            timing = sld.find(qn("p:timing"))
            if timing is not None:
                timing.addprevious(trans)
            else:
                sld.append(trans)
        except Exception:
            pass  # 动画注入失败不影响生成

    def _apply_entrances(self, slide, shape_ids: list[int], effect: str = "fade",
                         delay_step: int = 120):
        """为整页形状注入入场动画（单个 <p:timing> 内含多个效果）。"""
        pid = ENTRANCE_PRESETS.get(effect, 10)
        if pid == 0 or not shape_ids:
            return
        try:
            from lxml import etree

            ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

            parts = []
            for i, spid in enumerate(shape_ids):
                delay = i * delay_step
                parts.append(f"""
<par>
  <cTn id="3" fill="hold">
    <stCondLst><cond delay="indefinite"/></stCondLst>
    <childTnLst>
      <par>
        <cTn id="4" fill="hold">
          <stCondLst><cond delay="0"/></stCondLst>
          <childTnLst>
            <par>
              <cTn id="5" presetID="{pid}" presetClass="entr" presetSubtype="0" fill="hold" grpId="{i}" nodeType="clickEffect">
                <stCondLst><cond delay="{delay}"/></stCondLst>
                <childTnLst>
                  <set>
                    <cBhvr>
                      <cTn id="6" dur="1" fill="hold">
                        <stCondLst><cond delay="0"/></stCondLst>
                      </cTn>
                      <tgtEl><spTgt spid="{spid}"/></tgtEl>
                      <attrNameLst><attrName>style.visibility</attrName></attrNameLst>
                    </cBhvr>
                    <to><strVal val="visible"/></to>
                  </set>
                  <animEffect transition="in" filter="fade">
                    <cBhvr>
                      <cTn id="7" dur="500"/>
                      <tgtEl><spTgt spid="{spid}"/></tgtEl>
                    </cBhvr>
                  </animEffect>
                </childTnLst>
              </cTn>
            </par>
          </childTnLst>
        </cTn>
      </par>
    </childTnLst>
  </cTn>
</par>""")

            xml = f"""<p:timing xmlns:p="{ns_p}" xmlns:a="{ns_a}" xmlns:r="{ns_r}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                {''.join(parts)}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""

            el = etree.fromstring(xml.encode("utf-8"))
            sld = slide._element
            for old in sld.findall(qn("p:timing")):
                sld.remove(old)
            sld.append(el)
        except Exception:
            pass  # 动画注入失败不影响生成

    # ---------- 版式 ----------

    def _page_number(self, slide, n, total):
        self._textbox(slide, W - 1.6, H - 0.55, 1.3, 0.4,
                      f"{n:02d} / {total:02d}", size=11, color=self.theme.muted,
                      align=PP_ALIGN.RIGHT)

    def _accent_bar(self, slide, x, y, w=0.09, h=0.5, color=None):
        self._shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=color or self.accent)

    def _title_area(self, slide, text, y=0.55, size=34, page=None, total=None):
        self._accent_bar(slide, 0.75, y + 0.06, 0.09, 0.5)
        self._textbox(slide, 1.0, y, 10.8, 0.75, text, size=size,
                      color=self.theme.ink, bold=True, font=self.theme.heading_font)
        if page and total:
            self._page_number(slide, page, total)

    def _notes(self, slide, notes: str):
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # ---------- 各版式 ----------

    def _apply_accent(self, acc: Optional[str]):
        self._cur_accent = acc or self.theme.accent

    def build(self, outline: dict, effect: str = "fade",
              transition: str = "fade", animate: bool = True) -> io.BytesIO:
        """渲染 Outline 为 pptx 字节流。"""
        # 统一经由大纲校验/清洗，保证结构稳定
        outline = parse_outline(outline)
        self.effect = effect if effect in ENTRANCE_PRESETS else "fade"
        self.transition = transition
        total = len(outline["slides"])
        title = outline.get("title", "演示文稿")
        subtitle = outline.get("subtitle", "")
        author = outline.get("author", "")
        date = outline.get("date", "")
        theme = get_theme(outline.get("theme") or self.theme.id)
        self.theme = theme

        for i, sl in enumerate(outline["slides"]):
            stype = sl.get("type", "bullets")
            page = i + 1
            accent = sl.get("accent") or self.theme.accent
            self._apply_accent(accent)
            s = self._new_slide(accent)
            shapes = []

            if stype == "title":
                self._build_title(s, sl, title, subtitle, author, date)
            elif stype == "section":
                self._build_section(s, sl)
            elif stype == "two_column":
                self._build_two_column(s, sl, page, total)
            elif stype == "quote":
                self._build_quote(s, sl, page, total)
            elif stype == "stats":
                self._build_stats(s, sl, page, total)
            elif stype == "chart":
                self._build_chart(s, sl, page, total)
            elif stype == "closing":
                self._build_closing(s, sl)
            else:  # bullets 默认
                self._build_bullets(s, sl, page, total)

            self._notes(s, sl.get("notes", ""))

            if animate and self.effect != "none":
                ids = [shp.shape_id for shp in s.shapes]
                self._apply_entrances(s, ids, self.effect)

        out = io.BytesIO()
        self.prs.save(out)
        out.seek(0)
        return out

    # --- title ---
    def _build_title(self, s, sl, title, subtitle, author, date):
        accent = self.accent
        self._shape(s, MSO_SHAPE.OVAL, 10.2, -1.4, 3.6, 3.6, fill=accent)
        self._shape(s, MSO_SHAPE.OVAL, 10.9, 5.6, 3.0, 3.0, fill=None, line=accent)
        self._shape(s, MSO_SHAPE.OVAL, 0.6, 5.9, 1.0, 1.0, fill=self.theme.accent2)
        self._shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 2.62, 0.62, 0.12, fill=accent)
        self._textbox(s, 1.0, 2.85, 10.9, 1.7, title, size=46,
                      color=self.theme.ink, bold=True, font=self.theme.heading_font)
        sub = subtitle or sl.get("subtitle", "")
        if sub:
            self._textbox(s, 1.0, 4.55, 10.6, 0.7, sub, size=20,
                          color=self.theme.muted, font=self.theme.body_font)
        meta = " · ".join([x for x in (author, date) if x])
        if meta:
            self._textbox(s, 1.0, 6.35, 9.0, 0.5, meta, size=13,
                          color=self.theme.muted, font=self.theme.body_font)

    # --- section ---
    def _build_section(self, s, sl):
        accent = self.accent
        self._shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 1.35, 2.4, 1.05,
                    fill=accent, line=None)
        self._set_text_in_shape(s.shapes[-1], sl.get("title", "")[:12].replace("·", ""),
                                size=24, color=self.theme.bg if self.theme.is_dark else "#FFFFFF",
                                bold=True)
        self._textbox(s, 1.0, 3.0, 11.2, 1.2, sl.get("title", ""), size=40,
                      color=self.theme.ink, bold=True, font=self.theme.heading_font)
        bullets = sl.get("bullets") or sl.get("points") or []
        if bullets:
            self._bullets_frame(s, 1.05, 4.35, 10.8, 2.0, bullets, size=18,
                                color=self.theme.muted)

    # --- bullets ---
    def _build_bullets(self, s, sl, page, total):
        self._title_area(s, sl.get("title", ""), page=page, total=total)
        bullets = sl.get("bullets") or sl.get("points") or []
        self._bullets_frame(s, 1.05, 1.9, 11.2, 4.6, bullets, size=19, spacing=1.3)

    # --- two column ---
    def _build_two_column(self, s, sl, page, total):
        self._title_area(s, sl.get("title", ""), page=page, total=total)
        cols = sl.get("columns") or []
        if not cols:
            cols = [{"heading": "左", "bullets": []}, {"heading": "右", "bullets": []}]
        n = len(cols)
        gap, margin, top = 0.4, 1.0, 1.9
        cw = (W - 2 * margin - (n - 1) * gap) / n
        ch = 4.6
        for idx, col in enumerate(cols):
            x = margin + idx * (cw + gap)
            self._shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, cw, ch,
                        fill=self.theme.card, line=None)
            self._shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.25, top + 0.32, 0.5, 0.09,
                        fill=self.accent)
            self._textbox(s, x + 0.25, top + 0.55, cw - 0.5, 0.6,
                          col.get("heading", "要点"), size=20,
                          color=self.theme.ink, bold=True, font=self.theme.heading_font)
            self._bullets_frame(s, x + 0.3, top + 1.25, cw - 0.6, ch - 1.4,
                                col.get("bullets", []), size=15, color=self.theme.ink)

    # --- quote ---
    def _build_quote(self, s, sl, page, total):
        self._title_area(s, sl.get("title", "") or "引言", page=page, total=total)
        self._textbox(s, 1.0, 2.4, 11.2, 2.6, "\u201C" + sl.get("quote", ""),
                      size=30, color=self.theme.ink, bold=True,
                      font=self.theme.heading_font, line_spacing=1.35)
        src = sl.get("quote_source", "")
        if src:
            self._textbox(s, 1.05, 5.15, 11.0, 0.6, "— " + src, size=16,
                          color=self.accent, font=self.theme.body_font)

    # --- stats ---
    def _build_stats(self, s, sl, page, total):
        self._title_area(s, sl.get("title", ""), page=page, total=total)
        stats = sl.get("stats") or []
        if not stats:
            stats = [{"value": "—", "label": ""}]
        n = len(stats)
        gap, margin, top = 0.4, 1.0, 2.3
        cw = (W - 2 * margin - (n - 1) * gap) / n
        ch = 3.0
        for idx, st in enumerate(stats):
            x = margin + idx * (cw + gap)
            self._shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, cw, ch,
                        fill=self.theme.card, line=None)
            self._shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.3, top + 0.45, 1.1, 0.12,
                        fill=self.accent)
            self._textbox(s, x + 0.3, top + 0.85, cw - 0.6, 1.2,
                          st.get("value", ""), size=40, color=self.accent,
                          bold=True, font=self.theme.heading_font)
            self._textbox(s, x + 0.3, top + 2.15, cw - 0.6, 0.6,
                          st.get("label", ""), size=14, color=self.theme.muted,
                          font=self.theme.body_font)

    # --- chart ---
    def _build_chart(self, s, sl, page, total):
        self._title_area(s, sl.get("title", ""), page=page, total=total)
        chart = sl.get("chart") or {}
        ctype = chart.get("type", "bar")
        labels = chart.get("labels") or []
        datasets = chart.get("datasets") or []
        if not datasets:
            datasets = [{"name": "数据", "values": []}]

        cd = CategoryChartData()
        if labels:
            cd.categories = labels
        for d in datasets:
            vals = [v if isinstance(v, (int, float)) else 0 for v in (d.get("values") or [])]
            cd.add_series(d.get("name", "数据"), vals)

        x, y, w, h = 1.0, 1.9, 11.2, 4.7
        try:
            chart_type = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE_MARKERS,
                "pie": XL_CHART_TYPE.PIE,
                "donut": XL_CHART_TYPE.DOUGHNUT,
            }.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
            gf = s.shapes.add_chart(chart_type, Inches(x), Inches(y),
                                    Inches(w), Inches(h), cd)
            ch = gf.chart
            if ctype in ("bar", "line"):
                ch.has_legend = True
                ch.legend.position = XL_LEGEND_POSITION.BOTTOM
                ch.legend.include_in_layout = False
            # 主题化图表颜色
            try:
                from pptx.oxml.ns import nsmap  # noqa
            except Exception:
                pass
            self._tint_chart(ch, datasets)
        except Exception:
            # 图表失败时降级为数据表格文本
            lines = [f"{d.get('name','')}: " + ", ".join(str(v) for v in (d.get('values') or []))
                     for d in datasets]
            self._textbox(s, 1.0, 2.2, 11.2, 3.8, "\n".join(lines), size=16,
                          color=self.theme.ink, font=self.theme.body_font)

    def _tint_chart(self, chart, datasets):
        """为系列上色（尽力而为，失败无碍）。"""
        try:
            from pptx.oxml.ns import qn as _qn
            ser_list = chart._chartSpace.findall(".//" + _qn("c:ser"))
            palette = [self.theme.accent, self.theme.accent2,
                       "#FBBF24", "#A78BFA", "#34D399"]
            for i, ser in enumerate(ser_list):
                color = palette[i % len(palette)]
                # 替换 solidFill
                for sf in ser.findall(_qn("c:spPr") + "/" + _qn("a:solidFill")):
                    ser.find(_qn("c:spPr")).remove(sf)
                spPr = ser.find(_qn("c:spPr"))
                if spPr is None:
                    continue
                from lxml import etree
                solid = etree.SubElement(spPr, _qn("a:solidFill"))
                srgb = etree.SubElement(solid, _qn("a:srgbClr"))
                srgb.set("val", color.lstrip("#"))
        except Exception:
            pass

    # --- closing ---
    def _build_closing(self, s, sl):
        accent = self.accent
        self._shape(s, MSO_SHAPE.OVAL, -1.0, 5.3, 3.2, 3.2, fill=accent)
        self._shape(s, MSO_SHAPE.OVAL, 11.4, -1.0, 2.6, 2.6, fill=None, line=accent)
        self._textbox(s, 1.0, 2.7, 11.2, 1.4, sl.get("title", "谢谢"), size=52,
                      color=self.theme.ink, bold=True, font=self.theme.heading_font)
        sub = sl.get("subtitle", "")
        if sub:
            self._textbox(s, 1.0, 4.3, 11.2, 0.7, sub, size=18,
                          color=self.theme.muted, font=self.theme.body_font)
        self._textbox(s, 1.0, 6.4, 11.0, 0.5, "Generated with SlideForge", size=12,
                      color=self.theme.muted, font=self.theme.body_font,
                      align=PP_ALIGN.LEFT)


def build_pptx(outline: dict, effect: str = "fade",
               transition: str = "fade", animate: bool = True) -> io.BytesIO:
    """渲染 Outline 为 pptx 字节流（便捷入口）。"""
    builder = SlideBuilder(get_theme(outline.get("theme") or "ocean"),
                           effect=effect, transition=transition)
    return builder.build(outline, effect=effect, transition=transition, animate=animate)


def build_sample_pptx(theme: str = "ocean", effect: str = "fade",
                      transition: str = "fade") -> io.BytesIO:
    outline = sample_outline()
    outline["theme"] = theme
    return build_pptx(outline, effect=effect, transition=transition)
