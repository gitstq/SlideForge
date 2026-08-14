import io

from pptx import Presentation

from slideforge.builder import build_pptx, build_sample_pptx
from slideforge.outline import sample_outline


def _load(data: bytes) -> Presentation:
    return Presentation(io.BytesIO(data))


def test_sample_pptx_generates():
    data = build_sample_pptx(theme="ocean").getvalue()
    assert data[:4] == b"PK\x03\x04"  # zip magic
    prs = _load(data)
    assert len(prs.slides) == len(sample_outline()["slides"])


def test_all_themes_generate():
    from slideforge.themes import THEMES

    for tid in THEMES:
        data = build_sample_pptx(theme=tid).getvalue()
        assert _load(data)  # 可重新打开即合法


def test_16x9_dimensions():
    prs = _load(build_sample_pptx().getvalue())
    ratio = int(prs.slide_width) / int(prs.slide_height)
    assert abs(ratio - 16 / 9) < 0.01


def test_notes_present():
    outline = sample_outline()
    outline["slides"][2]["notes"] = "这是演讲备注"
    prs = _load(build_pptx(outline).getvalue())
    slides = list(prs.slides)
    note = slides[2].notes_slide.notes_text_frame.text
    assert "这是演讲备注" in note


def test_chart_present():
    prs = _load(build_sample_pptx().getvalue())
    any_chart = False
    for slide in prs.slides:
        if slide.shapes and any(s.has_chart for s in slide.shapes if hasattr(s, "has_chart")):
            any_chart = True
    assert any_chart


def test_empty_outline_rejected():
    import pytest

    from slideforge.outline import OutlineError

    with pytest.raises(OutlineError):
        build_pptx({"title": "x", "slides": []})
